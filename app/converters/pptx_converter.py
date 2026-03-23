from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from app.utils.file_utils import make_output_path


def pptx_to_txt(input_path: str) -> str:
    """Extract all text from every slide."""
    output_path = make_output_path(input_path, "txt")
    prs = Presentation(input_path)
    lines = []
    for i, slide in enumerate(prs.slides):
        lines.append(f"--- Slide {i + 1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
    return output_path


def pptx_to_pdf(input_path: str) -> str:
    """Convert PPTX slide text to PDF."""
    output_path = make_output_path(input_path, "pdf")
    prs = Presentation(input_path)
    styles = getSampleStyleSheet()
    story = []
    for i, slide in enumerate(prs.slides):
        story.append(Paragraph(f"Slide {i + 1}", styles["Heading2"]))
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        story.append(Paragraph(text, styles["Normal"]))
                        story.append(Spacer(1, 4))
    doc = SimpleDocTemplate(output_path, pagesize=letter)
    doc.build(story)
    return output_path


def pptx_to_html(input_path: str) -> str:
    """Convert PPTX slides to HTML."""
    output_path = make_output_path(input_path, "html")
    prs = Presentation(input_path)
    slides_html = ""
    for i, slide in enumerate(prs.slides):
        slides_html += f"<h2>Slide {i + 1}</h2>\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slides_html += f"<p>{text}</p>\n"
    html = f"""<!DOCTYPE html>
<html>
<head>
  <meta charset='utf-8'>
  <style>
    body {{ font-family: Arial, sans-serif; padding: 20px; }}
    h2 {{ color: #1A1A2E; border-bottom: 2px solid #0F3460; }}
  </style>
</head>
<body>
{slides_html}
</body>
</html>"""
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(html)
    return output_path